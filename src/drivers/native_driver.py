from pathlib import Path

from core.document_environment import BaseDocumentDriver, DocumentRequest, DocumentResult


class NativeDriver(BaseDocumentDriver):
    def supports(self, request: DocumentRequest) -> bool:
        return request.doc_type in ("docx", "xlsx", "pptx")

    def execute(self, request: DocumentRequest) -> DocumentResult:
        if request.action == "create":
            return self._create(request)
        return DocumentResult(success=False, message=f"Unsupported action: {request.action}")

    def _create(self, request: DocumentRequest) -> DocumentResult:
        output_path = Path(request.output_path) if request.output_path else None
        if not output_path:
            return DocumentResult(success=False, message="output_path required")
        output_path.parent.mkdir(parents=True, exist_ok=True)

        if request.doc_type == "docx":
            from docx import Document

            doc = Document()
            title = request.parameters.get("title", "")
            if title:
                doc.add_heading(title, level=1)
            doc.save(str(output_path))
        elif request.doc_type == "xlsx":
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws.title = "Sheet1"
            ws["A1"] = request.parameters.get("title", "")
            wb.save(str(output_path))
        elif request.doc_type == "pptx":
            from pptx import Presentation

            prs = Presentation()
            title = request.parameters.get("title", "")
            if title:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title
            prs.save(str(output_path))
        else:
            return DocumentResult(
                success=False, message=f"Unsupported doc_type: {request.doc_type}"
            )

        return DocumentResult(success=True, output_path=str(output_path), message="Created")

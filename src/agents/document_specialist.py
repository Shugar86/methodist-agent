"""
Document Specialist — агент для создания и редактирования документов.
Windows-first. Поддерживает COM-режим (pywin32) и Native-режим (fallback).
"""

import logging
from pathlib import Path
from typing import Any, Dict, Optional

from core.config import Config
from core.ui_text import (
    error_no_office_fallback,
    error_template_folder_empty,
    error_template_not_found,
    error_unsupported_format,
)
from .utils import ensure_directory, find_replace_docx, generate_filename

logger = logging.getLogger(__name__)

# Пути по умолчанию
TEMPLATES_DIR = Path.home() / ".methodist-agent" / "templates"
OUTPUT_DIR = Path.home() / ".methodist-agent" / "output"


class DocumentSpecialist:
    """
    Агент для работы с документами Word, Excel, PowerPoint.
    Приоритет: COM-режим (через pywin32). Fallback: native библиотеки.
    """

    def __init__(self, config: Optional[Config] = None):
        self.config = config
        self._word_app: Optional[Any] = None
        self._excel_app: Optional[Any] = None
        self._ppt_app: Optional[Any] = None
        self._com_available = self._check_com()
        self._native_available = self._check_native()
        logger.info(f"DocumentSpecialist init: COM={self._com_available}, Native={self._native_available}")

    # ------------------------------------------------------------------
    # Проверка доступности режимов
    # ------------------------------------------------------------------

    @staticmethod
    def _check_com() -> bool:
        try:
            import win32com.client  # noqa: F401
            return True
        except Exception:
            return False

    @staticmethod
    def _check_native() -> bool:
        try:
            import docx  # noqa: F401
            import openpyxl  # noqa: F401
            import pptx  # noqa: F401
            return True
        except Exception:
            return False

    # ------------------------------------------------------------------
    # Получение COM-объектов (с кэшированием)
    # ------------------------------------------------------------------

    def _get_word_app(self) -> Any:
        if not self._com_available:
            raise RuntimeError("COM недоступен")
        if self._word_app is None:
            import win32com.client as win32
            self._word_app = win32.Dispatch("Word.Application")
            self._word_app.Visible = True
            logger.info("Word COM запущен")
        return self._word_app

    def _get_excel_app(self) -> Any:
        if not self._com_available:
            raise RuntimeError("COM недоступен")
        if self._excel_app is None:
            import win32com.client as win32
            self._excel_app = win32.Dispatch("Excel.Application")
            self._excel_app.Visible = True
            logger.info("Excel COM запущен")
        return self._excel_app

    def _get_powerpoint_app(self) -> Any:
        if not self._com_available:
            raise RuntimeError("COM недоступен")
        if self._ppt_app is None:
            import win32com.client as win32
            self._ppt_app = win32.Dispatch("PowerPoint.Application")
            self._ppt_app.Visible = True
            logger.info("PowerPoint COM запущен")
        return self._ppt_app

    # ------------------------------------------------------------------
    # Диспетчер задач
    # ------------------------------------------------------------------

    def execute(self, task_type: str, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Главный метод диспетчеризации.

        task_type:
            create_docx, edit_docx,
            create_xlsx, edit_xlsx,
            create_pptx, edit_pptx,
            create_from_template
        """
        task_type = task_type.lower().strip()
        handlers = {
            "create_docx": self.create_docx,
            "edit_docx": self.edit_docx,
            "create_xlsx": self.create_xlsx,
            "edit_xlsx": self.edit_xlsx,
            "create_pptx": self.create_pptx,
            "edit_pptx": self.edit_pptx,
            "create_from_template": self.create_from_template,
        }

        handler = handlers.get(task_type)
        if handler is None:
            return {"success": False, "error": f"Неизвестный тип задачи: {task_type}"}

        try:
            return handler(params)
        except Exception as exc:
            logger.exception(f"Ошибка при выполнении {task_type}")
            return {"success": False, "error": str(exc)}

    # ------------------------------------------------------------------
    # Word
    # ------------------------------------------------------------------

    def create_docx(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Создание Word-документа.
        params:
            - content: список строк/абзацев
            - filename: имя файла (опционально)
            - subject: тема (для генерации имени)
            - template: путь к шаблону (опционально)
        """
        content = params.get("content", [])
        subject = params.get("subject", "документ")
        filename = params.get("filename") or generate_filename("doc", subject, "docx")
        template = params.get("template")

        output_path = ensure_directory(OUTPUT_DIR) / filename

        # Попытка COM
        if self._com_available:
            try:
                word = self._get_word_app()
                doc = word.Documents.Add()

                if template and Path(template).exists():
                    # В COM нельзя напрямую "применить" шаблон к существующему документу,
                    # но можно создать на основе шаблона
                    doc.Close(SaveChanges=False)
                    doc = word.Documents.Add(str(Path(template).resolve()))

                # Добавление текста
                for i, paragraph in enumerate(content):
                    if i == 0:
                        doc.Paragraphs(1).Range.Text = str(paragraph)
                    else:
                        doc.Content.InsertAfter(f"\n{paragraph}")

                doc.SaveAs(str(output_path.resolve()))
                logger.info(f"DOCX создан (COM): {output_path}")
                return {"success": True, "path": str(output_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM create_docx не удался: {exc}. Fallback на native.")

        # Fallback native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from docx import Document

        if template and Path(template).exists():
            doc = Document(template)
        else:
            doc = Document()

        for paragraph in content:
            doc.add_paragraph(str(paragraph))

        doc.save(str(output_path))
        logger.info(f"DOCX создан (native): {output_path}")
        return {"success": True, "path": str(output_path), "mode": "native"}

    def edit_docx(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Редактирование Word-документа.
        params:
            - path: путь к файлу
            - replacements: dict {что_искать: чем_заменить}
            - append: список строк для добавления в конец
        """
        file_path = Path(params.get("path", "")).expanduser()
        replacements = params.get("replacements", {})
        append = params.get("append", [])

        if not file_path.exists():
            return {"success": False, "error": f"Файл не найден: {file_path}"}

        # Попытка COM
        if self._com_available:
            try:
                word = self._get_word_app()
                doc = word.Documents.Open(str(file_path.resolve()))

                # Поиск и замена через COM
                for find_text, replace_text in replacements.items():
                    find = word.Selection.Find
                    find.Text = find_text
                    find.Replacement.Text = replace_text
                    find.Forward = True
                    find.Wrap = 1  # wdFindContinue
                    find.Format = False
                    find.MatchCase = False
                    find.MatchWholeWord = False
                    find.Execute(Replace=2)  # wdReplaceAll

                # Добавление текста
                if append:
                    end_range = doc.Content
                    end_range.Collapse(Direction=0)  # wdCollapseEnd
                    end_range.InsertAfter("\n" + "\n".join(str(a) for a in append))

                doc.Save()
                logger.info(f"DOCX отредактирован (COM): {file_path}")
                return {"success": True, "path": str(file_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM edit_docx не удался: {exc}. Fallback на native.")

        # Fallback native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from docx import Document

        doc = Document(str(file_path))

        for find_text, replace_text in replacements.items():
            find_replace_docx(doc, find_text, replace_text)

        for paragraph in append:
            doc.add_paragraph(str(paragraph))

        doc.save(str(file_path))
        logger.info(f"DOCX отредактирован (native): {file_path}")
        return {"success": True, "path": str(file_path), "mode": "native"}

    # ------------------------------------------------------------------
    # Excel
    # ------------------------------------------------------------------

    def create_xlsx(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Создание Excel-файла.
        params:
            - sheets: dict {имя_листа: [[строки], ...]}
            - filename: имя файла (опционально)
            - subject: тема
        """
        sheets = params.get("sheets", {"Лист1": []})
        subject = params.get("subject", "документ")
        filename = params.get("filename") or generate_filename("xls", subject, "xlsx")
        output_path = ensure_directory(OUTPUT_DIR) / filename

        # Попытка COM
        if self._com_available:
            try:
                excel = self._get_excel_app()
                wb = excel.Workbooks.Add()

                # Удаляем лишние листы, оставляем один
                while wb.Sheets.Count > 1:
                    wb.Sheets(1).Delete()

                for sheet_name, rows in sheets.items():
                    try:
                        ws = wb.Sheets.Add(After=wb.Sheets(wb.Sheets.Count))
                    except Exception:
                        ws = wb.Sheets(wb.Sheets.Count)
                    ws.Name = sheet_name
                    for r_idx, row in enumerate(rows, start=1):
                        for c_idx, value in enumerate(row, start=1):
                            ws.Cells(r_idx, c_idx).Value = value

                # Удаляем первый пустой лист, если создали новые
                if wb.Sheets.Count > 1:
                    try:
                        wb.Sheets(1).Delete()
                    except Exception:
                        pass

                wb.SaveAs(str(output_path.resolve()))
                logger.info(f"XLSX создан (COM): {output_path}")
                return {"success": True, "path": str(output_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM create_xlsx не удался: {exc}. Fallback на native.")

        # Fallback native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from openpyxl import Workbook

        wb = Workbook()
        first = True
        for sheet_name, rows in sheets.items():
            if first:
                ws = wb.active
                ws.title = sheet_name
                first = False
            else:
                ws = wb.create_sheet(title=sheet_name)
            for row in rows:
                ws.append(row)

        wb.save(str(output_path))
        logger.info(f"XLSX создан (native): {output_path}")
        return {"success": True, "path": str(output_path), "mode": "native"}

    def edit_xlsx(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Редактирование Excel-файла.
        params:
            - path: путь к файлу
            - cell_updates: dict {"Лист1!A1": значение, ...}
            - append_rows: dict {"Лист1": [[строки], ...]}
        """
        file_path = Path(params.get("path", "")).expanduser()
        cell_updates = params.get("cell_updates", {})
        append_rows = params.get("append_rows", {})

        if not file_path.exists():
            return {"success": False, "error": f"Файл не найден: {file_path}"}

        # Попытка COM
        if self._com_available:
            try:
                excel = self._get_excel_app()
                wb = excel.Workbooks.Open(str(file_path.resolve()))

                for addr, value in cell_updates.items():
                    if "!" in addr:
                        sheet_name, cell = addr.split("!", 1)
                        ws = wb.Sheets(sheet_name)
                    else:
                        ws = wb.ActiveSheet
                        cell = addr
                    ws.Range(cell).Value = value

                for sheet_name, rows in append_rows.items():
                    ws = wb.Sheets(sheet_name)
                    last_row = ws.UsedRange.Rows.Count
                    for r_idx, row in enumerate(rows, start=last_row + 1):
                        for c_idx, value in enumerate(row, start=1):
                            ws.Cells(r_idx, c_idx).Value = value

                wb.Save()
                logger.info(f"XLSX отредактирован (COM): {file_path}")
                return {"success": True, "path": str(file_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM edit_xlsx не удался: {exc}. Fallback на native.")

        # Fallback native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from openpyxl import load_workbook

        wb = load_workbook(str(file_path))

        for addr, value in cell_updates.items():
            if "!" in addr:
                sheet_name, cell = addr.split("!", 1)
                ws = wb[sheet_name]
            else:
                ws = wb.active
                cell = addr
            ws[cell] = value

        for sheet_name, rows in append_rows.items():
            ws = wb[sheet_name]
            for row in rows:
                ws.append(row)

        wb.save(str(file_path))
        logger.info(f"XLSX отредактирован (native): {file_path}")
        return {"success": True, "path": str(file_path), "mode": "native"}

    # ------------------------------------------------------------------
    # PowerPoint
    # ------------------------------------------------------------------

    def create_pptx(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Создание PowerPoint-презентации.
        params:
            - slides: список словарей:
                [{"title": "...", "content": ["..."], "layout": 1}, ...]
            - filename: имя файла (опционально)
            - subject: тема
        """
        slides = params.get("slides", [])
        subject = params.get("subject", "презентация")
        filename = params.get("filename") or generate_filename("ppt", subject, "pptx")
        output_path = ensure_directory(OUTPUT_DIR) / filename

        # Попытка COM
        if self._com_available:
            try:
                ppt = self._get_powerpoint_app()
                prs = ppt.Presentations.Add()

                for slide_data in slides:
                    title = slide_data.get("title", "")
                    content = slide_data.get("content", [])
                    layout_idx = slide_data.get("layout", 1)
                    slide = prs.Slides.Add(prs.Slides.Count + 1, layout_idx)
                    if slide.Shapes.Count >= 1:
                        slide.Shapes(1).TextFrame.TextRange.Text = title
                    if slide.Shapes.Count >= 2 and content:
                        slide.Shapes(2).TextFrame.TextRange.Text = "\n".join(str(c) for c in content)

                prs.SaveAs(str(output_path.resolve()))
                logger.info(f"PPTX создан (COM): {output_path}")
                return {"success": True, "path": str(output_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM create_pptx не удался: {exc}. Fallback на native.")

        # Fallback native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        for slide_data in slides:
            title = slide_data.get("title", "")
            content = slide_data.get("content", [])
            slide = prs.slides.add_slide(blank_layout)

            # Заголовок
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            # Контент
            top = Inches(1.7)
            height = Inches(5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content):
                if i == 0:
                    tf.text = str(line)
                else:
                    p = tf.add_paragraph()
                    p.text = str(line)
                    p.font.size = Pt(18)

        prs.save(str(output_path))
        logger.info(f"PPTX создан (native): {output_path}")
        return {"success": True, "path": str(output_path), "mode": "native"}

    def edit_pptx(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Редактирование PowerPoint-презентации.
        params:
            - path: путь к файлу
            - slide_updates: dict {slide_index: {"title": "...", "content": ["..."]}}
            - append_slides: список словарей (как в create_pptx)
        """
        file_path = Path(params.get("path", "")).expanduser()
        slide_updates = params.get("slide_updates", {})
        append_slides = params.get("append_slides", [])

        if not file_path.exists():
            return {"success": False, "error": f"Файл не найден: {file_path}"}

        # Попытка COM
        if self._com_available:
            try:
                ppt = self._get_powerpoint_app()
                prs = ppt.Presentations.Open(str(file_path.resolve()))

                for idx, update in slide_updates.items():
                    slide = prs.Slides(int(idx))
                    title = update.get("title")
                    content = update.get("content")
                    if title is not None and slide.Shapes.Count >= 1:
                        slide.Shapes(1).TextFrame.TextRange.Text = title
                    if content is not None and slide.Shapes.Count >= 2:
                        slide.Shapes(2).TextFrame.TextRange.Text = "\n".join(str(c) for c in content)

                for slide_data in append_slides:
                    title = slide_data.get("title", "")
                    content = slide_data.get("content", [])
                    layout_idx = slide_data.get("layout", 1)
                    slide = prs.Slides.Add(prs.Slides.Count + 1, layout_idx)
                    if slide.Shapes.Count >= 1:
                        slide.Shapes(1).TextFrame.TextRange.Text = title
                    if slide.Shapes.Count >= 2 and content:
                        slide.Shapes(2).TextFrame.TextRange.Text = "\n".join(str(c) for c in content)

                prs.Save()
                logger.info(f"PPTX отредактирован (COM): {file_path}")
                return {"success": True, "path": str(file_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM edit_pptx не удался: {exc}. Fallback на native.")

        # Fallback native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from pptx import Presentation

        prs = Presentation(str(file_path))

        for idx, update in slide_updates.items():
            slide = prs.slides[int(idx) - 1]
            title = update.get("title")
            content = update.get("content")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if title is not None and shape.text_frame.text.strip() == "":
                        shape.text_frame.text = title
                        title = None
                    elif content is not None:
                        shape.text_frame.text = "\n".join(str(c) for c in content)
                        content = None

        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        for slide_data in append_slides:
            title = slide_data.get("title", "")
            content = slide_data.get("content", [])
            slide = prs.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text_frame.text.strip() == "":
                        shape.text_frame.text = title
                        break

        prs.save(str(file_path))
        logger.info(f"PPTX отредактирован (native): {file_path}")
        return {"success": True, "path": str(file_path), "mode": "native"}

    # ------------------------------------------------------------------
    # Шаблоны
    # ------------------------------------------------------------------

    def create_from_template(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """
        Создание документа из шаблона с подстановкой.
        params:
            - template_name: имя файла шаблона (в ~/.methodist-agent/templates/)
            - subject: тема документа
            - replacements: dict {placeholder: значение}
            - output_format: docx | xlsx | pptx (по умолчанию из шаблона)
        """
        template_name = params.get("template_name", "")
        subject = params.get("subject", "документ")
        replacements = params.get("replacements", {})
        output_format = params.get("output_format")

        template_path = TEMPLATES_DIR / template_name
        if template_path.is_dir():
            candidates = ["template.docx", "template.xlsx", "template.pptx"]
            found = None
            for candidate in candidates:
                candidate_path = template_path / candidate
                if candidate_path.exists():
                    found = candidate_path
                    break
            if found is None:
                return {"success": False, "error": error_template_folder_empty(template_path)}
            template_path = found

        if not template_path.exists():
            return {"success": False, "error": error_template_not_found(template_path)}

        if output_format is None:
            output_format = template_path.suffix.lstrip(".").lower()

        filename = generate_filename(template_path.stem, subject, output_format)
        output_path = ensure_directory(OUTPUT_DIR) / filename

        # Word
        if output_format == "docx":
            return self._create_from_template_docx(template_path, output_path, replacements)

        # Excel
        if output_format == "xlsx":
            return self._create_from_template_xlsx(template_path, output_path, replacements)

        # PowerPoint
        if output_format == "pptx":
            return self._create_from_template_pptx(template_path, output_path, replacements)

        return {"success": False, "error": error_unsupported_format(output_format)}

    def _create_from_template_docx(self, template_path: Path, output_path: Path, replacements: Dict[str, str]) -> Dict[str, Any]:
        # COM
        if self._com_available:
            try:
                word = self._get_word_app()
                doc = word.Documents.Add(str(template_path.resolve()))
                for find_text, replace_text in replacements.items():
                    find = word.Selection.Find
                    find.Text = find_text
                    find.Replacement.Text = replace_text
                    find.Forward = True
                    find.Wrap = 1
                    find.Format = False
                    find.MatchCase = False
                    find.MatchWholeWord = False
                    find.Execute(Replace=2)
                doc.SaveAs(str(output_path.resolve()))
                return {"success": True, "path": str(output_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM create_from_template docx не удался: {exc}. Fallback на native.")

        # Native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from docx import Document
        doc = Document(str(template_path))
        for find_text, replace_text in replacements.items():
            find_replace_docx(doc, find_text, replace_text)
        doc.save(str(output_path))
        return {"success": True, "path": str(output_path), "mode": "native"}

    def _create_from_template_xlsx(self, template_path: Path, output_path: Path, replacements: Dict[str, str]) -> Dict[str, Any]:
        # COM
        if self._com_available:
            try:
                excel = self._get_excel_app()
                wb = excel.Workbooks.Open(str(template_path.resolve()))
                for sheet in wb.Sheets:
                    for addr, value in replacements.items():
                        try:
                            sheet.Range(addr).Value = value
                        except Exception:
                            pass
                wb.SaveAs(str(output_path.resolve()))
                return {"success": True, "path": str(output_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM create_from_template xlsx не удался: {exc}. Fallback на native.")

        # Native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from openpyxl import load_workbook
        wb = load_workbook(str(template_path))
        for addr, value in replacements.items():
            try:
                wb[addr] = value
            except Exception:
                pass
        wb.save(str(output_path))
        return {"success": True, "path": str(output_path), "mode": "native"}

    def _create_from_template_pptx(self, template_path: Path, output_path: Path, replacements: Dict[str, str]) -> Dict[str, Any]:
        # COM
        if self._com_available:
            try:
                ppt = self._get_powerpoint_app()
                prs = ppt.Presentations.Open(str(template_path.resolve()))
                for slide in prs.Slides:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            for find_text, replace_text in replacements.items():
                                if find_text in shape.TextFrame.TextRange.Text:
                                    shape.TextFrame.TextRange.Text = shape.TextFrame.TextRange.Text.replace(find_text, replace_text)
                prs.SaveAs(str(output_path.resolve()))
                return {"success": True, "path": str(output_path), "mode": "com"}
            except Exception as exc:
                logger.warning(f"COM create_from_template pptx не удался: {exc}. Fallback на native.")

        # Native
        if not self._native_available:
            return {"success": False, "error": error_no_office_fallback()}

        from pptx import Presentation
        prs = Presentation(str(template_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for find_text, replace_text in replacements.items():
                        if find_text in shape.text_frame.text:
                            shape.text_frame.text = shape.text_frame.text.replace(find_text, replace_text)
        prs.save(str(output_path))
        return {"success": True, "path": str(output_path), "mode": "native"}

    # ------------------------------------------------------------------
    # Очистка ресурсов
    # ------------------------------------------------------------------

    def close(self) -> None:
        """Закрывает COM-приложения."""
        if self._word_app:
            try:
                self._word_app.Quit()
            except Exception:
                pass
            self._word_app = None
        if self._excel_app:
            try:
                self._excel_app.Quit()
            except Exception:
                pass
            self._excel_app = None
        if self._ppt_app:
            try:
                self._ppt_app.Quit()
            except Exception:
                pass
            self._ppt_app = None
        logger.info("COM-приложения закрыты")
